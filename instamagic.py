#!/usr/bin/python3
import datetime
import sys
from itertools import dropwhile, takewhile
from os import listdir
from os.path import isfile, join

import instaloader
from pptx import Presentation
from pptx.util import Inches


def get_parameters():
    instagramid = input("Instagram ID: ")
    try:
        L = instaloader.Instaloader(save_metadata=False)
        instaloader.Profile.from_username(L.context, instagramid)
    except instaloader.exceptions.QueryReturnedNotFoundException:
        sys.exit(
            "Profile https://www.instagram.com/" +
            instagramid +
            " does not exist")

    print("Choose a start and end date. Format: 22-09-2020")
    from_date = input("From date: ")
    try:
        from_date = datetime.datetime.strptime(from_date, "%d-%m-%Y")
    except ValueError:
        sys.exit("invalid date format")

    end_date = input("End date: ")
    try:
        end_date = datetime.datetime.strptime(end_date, "%d-%m-%Y")
    except ValueError:
        sys.exit("invalid date format")

    print("Choose the order of images and videos in the presentation.")
    print("Choose 1 for newest to oldest, choose 2 for oldest to newest")
    order = input("Choose 1 or 2: ")
    if order != "1" and order != "2":
        sys.exit("Choose either 1 or 2")

    return instagramid, from_date, end_date, order


def download_posts(instagramid, from_date, end_date):
    L = instaloader.Instaloader(save_metadata=False)

    profile = instaloader.Profile.from_username(L.context, instagramid)
    posts = profile.get_posts()

    for post in takewhile(
        lambda p: p.date >= from_date,
        dropwhile(
            lambda p: p.date > end_date,
            posts)):
        print("Downloading: " + post.url)
        L.download_post(post, target=profile.username)


def pairwise(iterable):
    a = iter(iterable)
    return zip(a, a)


def make_presentation(instagramid, order):
    onlyfiles = [
        f for f in listdir(instagramid) if isfile(
            join(
                instagramid,
                f))]

    only_media_files = [f for f in onlyfiles if f.endswith(
        ".jpg") or f.endswith(".mp4")]

    only_media_files_no_poster_frames = only_media_files.copy()
    for file_name in only_media_files:
        if file_name.endswith(".mp4"):
            only_media_files_no_poster_frames.remove(
                file_name.replace(".mp4", ".jpg"))

    only_media_files_no_poster_frames.sort()
    if order == "1":
        only_media_files_no_poster_frames.reverse()

    prs = Presentation()
    layout6 = prs.slide_layouts[6]

    top = Inches(0.5)
    left = Inches(0.35)
    width = Inches(4.5)
    second = Inches(5.15)

    for file1, file2 in pairwise(only_media_files_no_poster_frames):
        slide = prs.slides.add_slide(layout6)

        if file1.endswith(".jpg"):
            slide.shapes.add_picture(
                instagramid + "/" + file1, left, top, width=width)
        if file1.endswith(".mp4"):
            slide.shapes.add_movie(
                instagramid +
                "/" +
                file1,
                left,
                top,
                width=width,
                height=width,
                poster_frame_image=instagramid +
                "/" +
                file1.replace(
                    ".mp4",
                    ".jpg"))
        if file2.endswith(".jpg"):
            slide.shapes.add_picture(
                instagramid + "/" + file2, second, top, width=width)
        if file2.endswith(".mp4"):
            slide.shapes.add_movie(
                instagramid +
                "/" +
                file2,
                second,
                top,
                width=width,
                height=width,
                poster_frame_image=instagramid +
                "/" +
                file2.replace(
                    ".mp4",
                    ".jpg"))

    prs.save(instagramid + "/" + instagramid + ".pptx")


def main():
    instagramid, from_date, end_date, order = get_parameters()

    download_posts(instagramid, from_date, end_date)

    make_presentation(instagramid, order)


if __name__ == "__main__":
    main()
